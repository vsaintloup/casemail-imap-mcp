from __future__ import annotations

from io import BytesIO
import csv
import logging
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader

from .models import Attachment
from .security import detect_prompt_injection_warnings

logger = logging.getLogger(__name__)


class AttachmentExtractionResult:
    def __init__(self, text: str | None, warnings: list[str] | None = None) -> None:
        self.text = text
        self.warnings = warnings or []


def _decode_text_bytes(payload: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return payload.decode(encoding)
        except UnicodeDecodeError:
            continue
    return payload.decode("utf-8", "replace")


def extract_text(
    attachment: Attachment,
    payload: bytes,
    max_bytes: int,
    max_chars: int,
    mode: str,
) -> AttachmentExtractionResult:
    if len(payload) > max_bytes:
        return AttachmentExtractionResult(None, ["Attachment exceeded configured extraction size limit."])

    filename = (attachment.filename or "").lower()
    mime_type = attachment.mime_type.lower()
    text: str | None = None
    warnings: list[str] = []

    try:
        if mime_type == "application/pdf" or filename.endswith(".pdf"):
            text = _extract_pdf(payload)
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or filename.endswith(".docx"):
            text = _extract_docx(payload)
        elif mime_type in {"text/plain", "text/markdown", "text/csv"} or filename.endswith((".txt", ".md", ".csv")):
            text = _extract_textual(payload, filename)
        elif filename.endswith(".xlsx"):
            text = _extract_xlsx(payload)
        elif filename.endswith(".pptx"):
            text = _extract_pptx(payload)
        elif mode == "ocr" or mime_type.startswith("image/"):
            text = _extract_ocr(payload)
            if text is None:
                warnings.append("OCR support is unavailable on this host or failed for this attachment.")
        elif mode == "all_small" and mime_type.startswith("text/"):
            text = _decode_text_bytes(payload)
    except Exception as exc:  # pragma: no cover - defensive parsing branch
        logger.warning("Attachment extraction failed: %s", str(exc))
        warnings.append("Attachment text extraction failed for this file type.")
        text = None

    if text:
        text = text.replace("\x00", "").strip()
        if len(text) > max_chars:
            text = text[:max_chars].rstrip() + "..."
        warnings.extend(detect_prompt_injection_warnings(text))
    return AttachmentExtractionResult(text=text, warnings=warnings)


def _extract_pdf(payload: bytes) -> str:
    reader = PdfReader(BytesIO(payload))
    return "\n".join((page.extract_text() or "").strip() for page in reader.pages if page.extract_text())


def _extract_docx(payload: bytes) -> str:
    document = DocxDocument(BytesIO(payload))
    return "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())


def _extract_textual(payload: bytes, filename: str) -> str:
    text = _decode_text_bytes(payload)
    if filename.endswith(".csv"):
        rows = list(csv.reader(text.splitlines()))
        return "\n".join(", ".join(cell for cell in row) for row in rows)
    return text


def _extract_xlsx(payload: bytes) -> str:
    workbook = load_workbook(BytesIO(payload), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in workbook.worksheets:
        lines.append(f"[Sheet] {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell) for cell in row if cell is not None and str(cell).strip()]
            if values:
                lines.append(" | ".join(values))
    return "\n".join(lines)


def _extract_pptx(payload: bytes) -> str:
    presentation = Presentation(BytesIO(payload))
    lines: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                slide_lines.append(text)
        if slide_lines:
            lines.append(f"[Slide {index}]")
            lines.extend(slide_lines)
    return "\n".join(lines)


def _extract_ocr(payload: bytes) -> str | None:
    try:
        import pytesseract
    except Exception:  # pragma: no cover - optional dependency runtime failure
        return None

    image = Image.open(BytesIO(payload))
    text = pytesseract.image_to_string(image)
    return text.strip() or None

