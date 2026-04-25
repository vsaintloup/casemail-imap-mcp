from __future__ import annotations

from datetime import UTC, datetime
from email.message import EmailMessage
from io import BytesIO
from pathlib import Path
import sys

import pytest
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "src"
if str(SRC) not in sys.path:
    sys.path.insert(0, str(SRC))

from casemail_imap_mcp.config import Settings  # noqa: E402


@pytest.fixture()
def settings(tmp_path: Path) -> Settings:
    cache_dir = tmp_path / "cache"
    return Settings(
        _env_file=None,
        imap_host="imap.test.local",
        imap_port=993,
        imap_username="lawyer@example.com",
        imap_password="secret-password",
        imap_use_ssl=True,
        case_folder_allowlist_regex=r"^Client/.+",
        sent_folder_allowlist_regex=r"^(Sent|Sent Items)$",
        default_sent_folders="Sent,Sent Items",
        message_ref_secret="super-secret-message-ref-key",
        cache_db_path=cache_dir / "casemail.sqlite3",
        cache_key_path=cache_dir / "casemail.key",
    )


def make_email(
    *,
    subject: str,
    sender: str,
    to: list[str],
    cc: list[str] | None = None,
    date_value: datetime | None = None,
    body_text: str = "",
    body_html: str | None = None,
) -> bytes:
    message = EmailMessage()
    message["Subject"] = subject
    message["From"] = sender
    message["To"] = ", ".join(to)
    if cc:
        message["Cc"] = ", ".join(cc)
    message["Date"] = (date_value or datetime(2026, 2, 2, 15, 30, tzinfo=UTC)).strftime("%a, %d %b %Y %H:%M:%S %z")
    message["Message-ID"] = f"<{abs(hash(subject + sender))}@example.test>"
    if body_html is not None:
        message.set_content(body_text or "Plain fallback")
        message.add_alternative(body_html, subtype="html")
    else:
        message.set_content(body_text)
    return message.as_bytes()


@pytest.fixture()
def sample_docx_bytes() -> bytes:
    document = DocxDocument()
    document.add_paragraph("DOCX evidence line")
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


@pytest.fixture()
def sample_xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Activity"
    sheet["A1"] = "Task"
    sheet["B1"] = "Hours"
    sheet["A2"] = "Draft motion"
    sheet["B2"] = 2.5
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


@pytest.fixture()
def sample_pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    textbox.text_frame.text = "Hearing prep checklist"
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
